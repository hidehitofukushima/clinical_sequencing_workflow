from pptx import Presentation
import datetime


prs = Presentation()
prs = Presentation('/Users/fukushimahideto/Dropbox/NGS annotation-module/master.pptx')



#------------------------------------------------------------------------------------------------------
title_slide_layout = prs.slide_layouts[0]        # "タイトルとスライド"のSlideLayoutオブジェクトを取得
sld0 = prs.slides.add_slide(title_slide_layout)  # 新規スライドを追加

for p in sld0.placeholders:  # スライド内の個別PlaceFoldeオブジェクトを展開
    print(p.name)  # フォルダ名を取得
    # >>Title 1, Subtitle 2
    
    print(p.placeholder_format.idx, p.placeholder_format.type)  # インデックスと種類を取得
    # >>0 CENTER_TITLE (3), 1 SUBTITLE (4)

#------------------------------------------------------------------------------------------------------
title_slide_layout = prs.slide_layouts[8]        # "2つコンテンツ"のSlideLayoutオブジェクトを取得
sld1 = prs.slides.add_slide(title_slide_layout)  # 新規スライドを追加  

for p in sld1.placeholders:  # スライド内の個別PlaceFoldeオブジェクトを展開
    print(p.name)  # フォルダ名を取得    
    # >>Title 1, Content Placeholder 2, Content Placeholder 3
    
    print(p.placeholder_format.idx, p.placeholder_format.type)  # インデックスと種類を取得
    # >>0 TITLE (1), 1 OBJECT (7), 2 OBJECT (7)

#------------------------------------------------------------------------------------------------------
title = sld0.placeholders[0]                     # スライド1のTitle1のPlaceFolderオブジェクトの取得
title.text = "Clinical Sequence Annotation" 
subtitle = sld0.placeholders[1]
subtitle.text = "hidehito fukushima" # テキストを設定

title = sld1.placeholders[0]
sld1.placeholders[1].insert_picture('output/igv_snv/squish_2_47702377.png')# スライド2のTitle1のPlaceFolderオブジェクトの取得
title.text = "「2つのコンテンツ」のスライド"       # テキストを設定
prs.save('/Users/fukushimahideto/Dropbox/NGS annotation-module/master.pptx')	# 名前をつけてpptxファイルとして保存



#ID${ID}.pptx
#