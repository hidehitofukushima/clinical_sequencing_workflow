# coding: UTF-8

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import datetime

# センチメートルをPowerPoint上の距離に変換する関数
def Centis(length):
    centi = Inches(length/2.54)
    return centi


# ---------------------------------------------------
"テンプレート、出力ファイル名の設定"
# ---------------------------------------------------
# templateとなるpptxファイルを指定する。
template_path = "./template.pptx"
# 出力するpptxファイルを指定する。(存在しない場合、自動作成されます)
save_path = "./output.pptx"


###########################################################################################
# 空のスライドの挿入
###########################################################################################
# ---------------------------------------------------
presentation = Presentation(template_path)
if len(presentation.slides) > 0:
    rId = presentation.slides._sldIdLst[0].rId
    presentation.part.drop_rel(rId)
    del presentation.slides._sldIdLst[0]



###########################################################################################
# タイトルスライド
###########################################################################################

title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
# ---------------------------------------------------
"タイトルテキストの挿入"
# ---------------------------------------------------
# 1
# タイトル
id = 765
title_slide.shapes.placeholders[0].text = f"NGS analysis 1st Report\nID {id}"
title_slide.shapes.placeholders[1].text = f"Curated_by: Hidehito Fukushima"
title_slide.shapes.placeholders[0].text_frame.paragraphs[0].font.size = Pt(40)
title_slide.shapes.placeholders[0].text_frame.paragraphs[1].font.size = Pt(80)
title_slide.shapes.placeholders[0].text_frame.paragraphs[1].font.color.rgb = RGBColor(245,66,66) # もっと綺麗な赤色

# 2 右下
# 右下に所属などの情報を記載
text_box = title_slide.shapes.add_textbox(Centis(27), Centis(13.4), Centis(10), Centis(5)) # left, top, width, height
text_box.text = "IMSUT Hospital\nHematology/Oncology\nHidehito Fukushima\nTomokazu Seki\nTakamori Hiroyuki\nKazuaki Yokoyama"
text_box.text_frame.add_paragraph().font.size = Pt(20)
# text_box.fill.solid()
# text_box.fill.fore_color.rgb = color

#3 右上
# 日時
text_box = title_slide.shapes.add_textbox(Centis(27), Centis(1), Centis(10), Centis(5)) # left, top, width, height
text_box.text = f"Date: {datetime.datetime.now().strftime('%Y-%m-%d')}"
text_box.text_frame.add_paragraph().font.name = "Meiryo"
text_box.text_frame.add_paragraph().font.size = Pt(20)





# # ---------------------------------------------------
# "imageの挿入"
# # ---------------------------------------------------
# # 挿入する位置
# pic_left = Centis(1)
# pic_top = Centis(5)

# # imageの高さを指定
# pic_height = Centis(7.9)

# image_path = "./graph.png"
# slide.shapes.add_picture(image_path, pic_left, pic_top, height=pic_height)


# # ---------------------------------------------------
# "tableの挿入"
# # ---------------------------------------------------
# # 入力したいtable状のデータ
# sample_table = [ ["1.1","1.2","1.3"]
#                 ,["2.1","2.2","2.3"]
#                 ,["3.1","3.2","3.3"]]

# # cell内のフォントサイズ
# cell_font = 20

# # 挿入する位置
# table_left = Centis(9.4)
# table_top = Centis(5)

# # tableの幅と高さ（仮）
# table_width = Centis(15)
# table_height = Centis(10)

# # tableの行数と列数(tableのサイズ)
# rows = len(sample_table)
# cols = len(sample_table[0])

# table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table

# # 表の各セルの中身を記入
# for i in range(rows):
#     for j in range(cols):
#         cell = table.cell(i, j)
#         cell.text = sample_table[i][j]
#         cell.text_frame.paragraphs[0].font.size = Pt(cell_font)

# # tableの高さを再調整
# table.rows[0].height = Centis(1.5)
# table.rows[1].height = Centis(4.9)
# table.rows[2].height = Centis(1.5)

# # tableの幅を再調整
# table.columns[0].width = Centis((15) / 3)
# table.columns[1].width = Centis((15) / 3)
# table.columns[2].width = Centis((15) / 3)


# # ---------------------------------------------------
# "テキストボックスの挿入"
# # ---------------------------------------------------
# # 文字列


# # ---------------------------------------------------
# "ファイル保存"
# # ---------------------------------------------------
presentation.save(save_path)
